# -*- encoding: utf-8 -*-

# @File        :   pptx_parser.py
# @Time        :   2021/04/13 15:37:44
# @Author      :   fanchunke
# @Email       :   fanchunke@laiye.com
# @Description :   

import logging
import os
from typing import List

import pptx

from .base import BaseParser

logger = logging.getLogger("doc-extracter")


class Parser(BaseParser):

    supported_extension = '.pptx'

    @classmethod
    def parse(cls, filename: str) -> List[dict]:
        """ 解析 .pptx

        Args:
            filename (str): `.pptx` 为扩展名的文件

        Raises:
            Exception: 不支持的文件类型/文件不存在

        Returns:
            List[dict]: 解析结果
        """

        if not os.path.exists(filename):
            raise Exception(f"Not Found: {filename}")

        text_runs = []
        with open(filename, 'rb') as f:
            presentation = pptx.Presentation(f)
            for index, slide in enumerate(presentation.slides):
                contents = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            contents.append(run.text)
                
                origin_contents = "\n\n".join(contents)
                text_runs.append(
                    {
                        "page": str(index + 1),
                        "context": origin_contents,
                    }
                )
        return text_runs
